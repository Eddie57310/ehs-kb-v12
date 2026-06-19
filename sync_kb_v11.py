import os, glob, logging, subprocess, re, shutil, base64
from datetime import datetime
from pathlib import Path
import requests
import pypdf
import docx2txt
from docx import Document as DocxDocument
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import MarkdownTextSplitter, RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from pdf_table_extractor import (
    is_table_garbled, pdf_table_to_documents, pymupdf_to_documents, engine_status,
)

# ===================== 1. 核心路径与配置 =====================
BASE_DIR = os.path.expanduser("~/doc_parser_v12")
KB_DIR = f"{BASE_DIR}/Local_KB"
OUT_DIR = f"{BASE_DIR}/batch_output"
DB_DIR = f"{BASE_DIR}/chroma_db"
LOG_DIR = f"{BASE_DIR}/logs"
FAILED_DIR    = f"{BASE_DIR}/Failed_PDFs"
SLIDE_IMG_DIR = f"{BASE_DIR}/slide_images"
WEAK_PAGE_THRESHOLD  = 50
VISION_TEXT_THRESHOLD = 100   # 低于此字数的页面走 Vision LLM
VISION_API_KEY  = "ark-a1563795-462d-4fd6-893a-83cec208a8e9-44895"
VISION_BASE_URL = "https://ark.cn-beijing.volces.com/api/coding/v3/chat/completions"
VISION_MODEL    = "doubao-seed-2.0-pro"

os.makedirs(OUT_DIR,       exist_ok=True)
os.makedirs(DB_DIR,        exist_ok=True)
os.makedirs(LOG_DIR,       exist_ok=True)
os.makedirs(FAILED_DIR,    exist_ok=True)
os.makedirs(SLIDE_IMG_DIR, exist_ok=True)

current_date = datetime.now().strftime('%Y-%m-%d')
log_file = os.path.join(LOG_DIR, f"sync_kb_v11_{current_date}.log")
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[logging.FileHandler(log_file, encoding='utf-8'), logging.StreamHandler()]
)
logger = logging.getLogger(__name__)

# ===================== 2. 组织架构检测与 SQLite 写入 =====================


def extract_docx_structured_chunks(file_path: str, max_size: int = 600) -> list[str]:
    """
    用 python-docx 按标题边界提取DOCX内容，每个chunk带章节路径前缀。
    表格输出为 Markdown 格式，与正文段落按阅读顺序交错处理。

    每个chunk格式：【一级标题 > 二级标题】\\n正文内容
    若章节内容超过 max_size，进一步按句子切分，但保留标题前缀。
    无标题时降级为 RecursiveCharacterTextSplitter。
    """
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph as _Para
    from docx.table import Table as _Table

    _sub = RecursiveCharacterTextSplitter(
        chunk_size=max_size, chunk_overlap=80,
        separators=["\n\n", "\n", "。", "；", "，", " "]
    )

    def _heading_level(style_name: str) -> int:
        for i in range(1, 7):
            if f'Heading {i}' in style_name or f'标题 {i}' in style_name or style_name == f'标题{i}':
                return i
        return 0

    def _table_to_md(table) -> str:
        """将 DOCX 表格转为 Markdown 表格字符串。"""
        rows_md = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            rows_md.append('| ' + ' | '.join(cells) + ' |')
            if i == 0:
                rows_md.append('|' + '|'.join(['---'] * len(cells)) + '|')
        return '\n'.join(rows_md)

    doc = DocxDocument(file_path)
    heading_stack: list[tuple[int, str]] = []
    body_lines: list[str] = []
    chunks: list[str] = []

    def _flush():
        content = '\n'.join(ln for ln in body_lines if ln.strip())
        if not content.strip():
            return
        prefix = ('【' + ' > '.join(h[1] for h in heading_stack) + '】\n') if heading_stack else ''
        full = prefix + content
        if len(full) <= max_size:
            chunks.append(full)
        else:
            for sub in _sub.split_text(content):
                chunks.append(prefix + sub)

    # 按阅读顺序遍历段落和表格（doc.paragraphs 会跳过表格）
    for child in doc.element.body:
        if child.tag == qn('w:p'):
            para = _Para(child, doc)
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ''
            level = _heading_level(style)
            if level > 0:
                _flush()
                body_lines = []
                while heading_stack and heading_stack[-1][0] >= level:
                    heading_stack.pop()
                heading_stack.append((level, text))
            else:
                body_lines.append(text)
        elif child.tag == qn('w:tbl'):
            table = _Table(child, doc)
            md_table = _table_to_md(table)
            if md_table.strip():
                body_lines.append(md_table)

    _flush()

    # 无标题文档降级为普通切片（保底）
    if not chunks:
        full_text = '\n'.join(p.text.strip() for p in doc.paragraphs if p.text.strip())
        chunks = _sub.split_text(full_text)

    return chunks



# ===================== 3. 工具函数 =====================
def extract_time_from_filename(filename, rel_path=None):
    """从文件名提取日期，找不到则向上查父目录，再找不到返回默认值。"""
    def _parse(s):
        m = re.search(r'(\d{4})[-_]?(\d{2})[-_]?(\d{2})', s)
        if m:
            year, month, day = m.groups()
            date_str = f"{year}-{month}-{day}"
            try:
                dt = datetime.strptime(date_str, "%Y-%m-%d")
                return int(dt.timestamp()), date_str
            except ValueError:
                pass
        return None

    result = _parse(filename)
    if result:
        return result

    # 文件名无日期时，逐级检查父目录
    if rel_path:
        parts = Path(rel_path).parent.parts
        for part in reversed(parts):  # 从最近父目录往上
            result = _parse(part)
            if result:
                return result

    return 946684800, "2000-01-01"

def convert_doc_to_docx(doc_path):
    out_dir = os.path.join(OUT_DIR, "_doc_converted")
    os.makedirs(out_dir, exist_ok=True)
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "docx", "--outdir", out_dir, doc_path],
            capture_output=True, timeout=300
        )
        base = os.path.splitext(os.path.basename(doc_path))[0]
        new_path = os.path.join(out_dir, base + ".docx")
        return new_path if os.path.exists(new_path) else None
    except Exception as e:
        logger.error(f"    ❌ LibreOffice 转换异常: {e}")
        return None

def convert_ods_to_xlsx(ods_path):
    out_dir = os.path.join(OUT_DIR, "_ods_converted")
    os.makedirs(out_dir, exist_ok=True)
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "xlsx", "--outdir", out_dir, ods_path],
            capture_output=True, timeout=300
        )
        base = os.path.splitext(os.path.basename(ods_path))[0]
        new_path = os.path.join(out_dir, base + ".xlsx")
        return new_path if os.path.exists(new_path) else None
    except Exception as e:
        logger.error(f"    ❌ LibreOffice ODS→XLSX 转换异常: {e}")
        return None

def find_weak_pages(md_path):
    try:
        with open(md_path, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception as e:
        logger.warning(f"    ⚠️ 无法读取MD进行体检: {e}")
        return []
    page_pattern = re.compile(r'\n---\s*PAGE\s+(\d+)\s*---\n', re.IGNORECASE)
    parts = page_pattern.split(content)
    weak_pages = []
    if parts and not page_pattern.match(parts[0]) and len(parts[0].strip()) < WEAK_PAGE_THRESHOLD:
        weak_pages.append(1)
    i = 1
    while i < len(parts) - 1:
        try:
            if len(parts[i + 1].strip()) < WEAK_PAGE_THRESHOLD:
                weak_pages.append(int(parts[i]))
            i += 2
        except Exception:
            i += 1
    return weak_pages

def extract_rescue_pdf(source_pdf_path, weak_pages):
    root_dir = os.path.dirname(source_pdf_path)
    file_name_no_ext = os.path.splitext(os.path.basename(source_pdf_path))[0]
    pages_str = "-".join(map(str, weak_pages))
    rescue_pdf_path = os.path.join(root_dir, f"{file_name_no_ext}_rescue_p{pages_str}.pdf")
    if os.path.exists(rescue_pdf_path):
        return rescue_pdf_path
    try:
        reader = pypdf.PdfReader(source_pdf_path)
        total_pages = len(reader.pages)
        writer = pypdf.PdfWriter()
        valid_pages = [p for p in weak_pages if 1 <= p <= total_pages]
        for p in valid_pages:
            writer.add_page(reader.pages[p - 1])
        if not valid_pages:
            return None
        with open(rescue_pdf_path, 'wb') as f:
            writer.write(f)
        logger.info(f"    🚑 救援PDF已生成，准备二次强行解析...")
        return rescue_pdf_path
    except Exception as e:
        logger.error(f"    ❌ 救援PDF提取失败: {e}")
        return None

def get_cached_md(output_dir, base_name):
    md_dir = os.path.join(output_dir, base_name)
    md_files = glob.glob(os.path.join(md_dir, "**", "*.md"), recursive=True)
    if md_files:
        return md_files[0]
    def clean_str(s): return re.sub(r'[^a-zA-Z0-9\u4e00-\u9fa5]', '', s)
    target_clean = clean_str(base_name)
    for root, _, files in os.walk(output_dir):
        for f in files:
            if f.endswith('.md') and clean_str(f[:-3]) == target_clean:
                logger.info(f"    🔍 模糊缓存命中: {f}")
                return os.path.join(root, f)
    return None

def run_mineru(pdf_path, output_dir):
    pdf_name = os.path.basename(pdf_path)
    base_name = os.path.splitext(pdf_name)[0]
    cached_md = get_cached_md(output_dir, base_name)
    if cached_md:
        logger.info(f"    ⚡ 发现解析缓存，直接读取！")
        return cached_md
    logger.info(f"    📕 MinerU 正在解析: {pdf_name}")
    try:
        res = subprocess.run(
            [os.path.join(BASE_DIR, "venv/bin/magic-pdf"), "-p", pdf_path, "-o", output_dir, "-m", "auto"],
            shell=False, capture_output=True, text=True, timeout=3600
        )
        if res.returncode != 0:
            logger.error(f"    ❌ MinerU 执行失败: {res.stderr[:200]}")
            return None
        return get_cached_md(output_dir, base_name)
    except subprocess.TimeoutExpired:
        logger.error(f"    ⏱️ MinerU 超时 (3600s): {pdf_name}")
        return None
    except Exception as e:
        logger.error(f"    ❌ MinerU 异常: {e}")
        return None

def move_to_failed_zone(file_path, file_name, reason):
    try:
        logger.error(f"    ☣️ 移入失败隔离区: {file_name}")
        shutil.copy(file_path, os.path.join(FAILED_DIR, file_name))
        with open(os.path.join(FAILED_DIR, f"{file_name}_error.txt"), "w") as f:
            f.write(f"失败原因: {reason}\n时间: {datetime.now()}")
    except Exception as e:
        logger.error(f"    ⚠️ 移入隔离区失败: {e}")

# ===================== 3. PDF 章节感知切块 =====================

def extract_pdf_structured_chunks(md_content: str, max_size: int = 800) -> list[str]:
    """
    对 MinerU 输出的 Markdown 按章节边界切块，每块带完整章节路径前缀。

    检测规则：
      1. Markdown 标题行：# 一级  ## 二级  ### 三级
      2. 数字编号式标题：1.  2.1  3.2.1（行尾为中文，不超过 40 字）
    每块格式：【一级 > 二级 > 三级】\\n正文内容
    超长块进一步分割，但保留前缀。
    无任何标题时降级为 RecursiveCharacterTextSplitter。
    """
    _sub = RecursiveCharacterTextSplitter(
        chunk_size=max_size, chunk_overlap=80,
        separators=["\n\n", "\n", "。", "；", "，", " "]
    )

    md_heading_re  = re.compile(r'^(#{1,6})\s+(.+)$')
    # 数字编号标题：裸数字(1 总则)与多级(5.1.5)都认；标题以句末标点(。；，、)结尾的当正文(列表项/句子)
    num_heading_re = re.compile(r'^(\d+(?:\.\d+)*)\s+([A-Za-z\u4e00-\u9fa5][\u4e00-\u9fa5\w\s（）【】、，。：:]{1,60})\s*$')

    def _detect_heading(line):
        m = md_heading_re.match(line)
        if m:
            return len(m.group(1)), m.group(2).strip()
        m = num_heading_re.match(line)
        if m:
            t = m.group(2).strip()
            if t and t[-1] in '。；，、':      # 句末标点结尾 → 列表项/句子，非标题
                return None
            num   = m.group(1)
            title = f"{num} {t}"
            level = num.count('.') + 1
            return level, title
        return None

    heading_stack: list[tuple[int, str]] = []
    body_lines: list[str] = []
    chunks: list[str] = []

    def _flush():
        content = '\n'.join(ln for ln in body_lines if ln.strip())
        if not content.strip():
            return
        prefix = ('【' + ' > '.join(h[1] for h in heading_stack) + '】\n') if heading_stack else ''
        full = prefix + content
        if len(full) <= max_size:
            chunks.append(full)
        else:
            for sub in _sub.split_text(content):
                chunks.append(prefix + sub)

    for line in md_content.splitlines():
        heading = _detect_heading(line.rstrip())
        if heading:
            _flush()
            body_lines = []
            level, title = heading
            while heading_stack and heading_stack[-1][0] >= level:
                heading_stack.pop()
            heading_stack.append((level, title))
        else:
            body_lines.append(line)
    _flush()

    if not chunks:
        return _sub.split_text(md_content)
    return chunks


# ===================== 3. OCR 噪音检测 & LLM 校对 =====================

def _is_toc_chunk(text: str) -> bool:
    """检测是否为目录页块（大量省略号+页码，无实质内容）。"""
    if not text or len(text) < 20:
        return False
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    if not lines:
        return False
    # 超过40%的行以数字结尾（页码）且含省略号
    toc_lines = sum(1 for l in lines if re.search(r'[.…]{3,}', l) and re.search(r'\d+\s*$', l))
    return toc_lines / len(lines) > 0.4


def _has_ocr_noise(text: str) -> bool:
    """检测文本是否含有明显的 OCR 噪音（随机 ASCII 串混入中文）。"""
    if not text:
        return False
    # 连续3个以上的 ASCII 非空格字符（数字/字母混排），且文本中有中文
    has_chinese = bool(re.search(r'[\u4e00-\u9fa5]', text))
    noise_segs = re.findall(r'[A-Za-z0-9]{3,}', text)
    # 过滤掉正常的英文单词（纯字母 ≤ 8 字符可能是合法词）
    noise_segs = [s for s in noise_segs if not s.isalpha() or len(s) > 8]
    return has_chinese and len(noise_segs) >= 2


OCR_BATCH_SIZE  = 20   # 每批校对块数，批间暂停避免限流
OCR_BATCH_PAUSE = 5    # 批间暂停秒数
OCR_CALL_DELAY  = 0.3  # 单次调用间隔秒数


def llm_correct_ocr(text: str) -> str:
    """
    用 LLM 清洗 OCR 噪音：删除随机字符串，修复断行，保留原始内容。
    只处理含噪音的块，跳过干净文本。
    """
    if not _has_ocr_noise(text):
        return text
    prompt = (
        "以下是通过OCR识别的中文PDF文本，可能含有识别错误（如混入随机英文字母串、"
        "数字噪音、断行错位等）。\n"
        "请保留所有实质性中文内容，删除明显的OCR噪音（随机英文字母串、无意义数字串），"
        "合并因扫描换行产生的断词，不要添加、删除或改变任何实质内容。"
        "直接输出修正后的文本，不要任何说明。\n\n"
        f"原文：\n{text}"
    )
    try:
        headers = {"Authorization": f"Bearer {VISION_API_KEY}", "Content-Type": "application/json"}
        payload = {
            "model": VISION_MODEL,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.1,
        }
        res = requests.post(VISION_BASE_URL, headers=headers, json=payload, timeout=60)
        data = res.json()
        if "choices" in data:
            return data["choices"][0]["message"]["content"]
    except Exception as e:
        logger.warning(f"    ⚠️ LLM校对失败，保留原文: {e}")
    return text


# ===================== 4. PDF 三引擎路由 =====================

def _inject_page_markers(md_content: str, md_path: str) -> str:
    """
    从 MinerU 生成的 content_list.json 中读取 page_idx，
    在 markdown 内容里插入 '--- PAGE N ---' 标记，供后续 chunk 页码追踪使用。
    MinerU 本身不在 .md 文件中写入页码标记，需要我们从 JSON 补充。
    """
    import json as _json
    md_dir = os.path.dirname(md_path)
    md_base = os.path.splitext(os.path.basename(md_path))[0]
    cl_path = os.path.join(md_dir, f"{md_base}_content_list.json")
    if not os.path.exists(cl_path):
        return md_content
    try:
        with open(cl_path, encoding='utf-8') as f:
            items = _json.load(f)
    except Exception:
        return md_content

    # 收集每页第一个非空文本（page_idx 从 0 开始，对应页码 = page_idx + 1）
    page_first_text: dict[int, str] = {}
    for item in items:
        p = item.get('page_idx')
        if p is None:
            continue
        text = item.get('text', '').strip()
        if p not in page_first_text and text:
            page_first_text[p] = text

    if not page_first_text:
        return md_content

    # 按页码升序扫描，记录每页标记应插入的位置
    insertions: list[tuple[int, int]] = []  # (pos_in_md, page_number)
    search_from = 0
    for page_idx in sorted(page_first_text.keys()):
        if page_idx == 0:
            continue  # 第一页不需要标记（默认 page=1）
        text = page_first_text[page_idx]
        page_num = page_idx + 1
        for n in (40, 25, 15, 8):
            snippet = text[:n].strip()
            if not snippet:
                continue
            pos = md_content.find(snippet, search_from)
            if pos != -1:
                insertions.append((pos, page_num))
                search_from = pos + 1
                break

    if not insertions:
        return md_content

    # 从后往前插入，避免位置偏移
    insertions.sort(key=lambda x: x[0], reverse=True)
    result = md_content
    for pos, page_num in insertions:
        result = result[:pos] + f'\n--- PAGE {page_num} ---\n' + result[pos:]
    return result


def _pdf_has_text_layer(pdf_path: str, min_chars_per_page: int = 50) -> bool:
    """pymupdf 快速探测 PDF 是否有文字层（非扫描件）。"""
    try:
        import fitz
        doc = fitz.open(pdf_path)
        total = sum(len(p.get_text()) for p in doc)
        avg = total / max(len(doc), 1)
        doc.close()
        return avg >= min_chars_per_page
    except Exception:
        return True  # 探测失败时保守假设有文字层


def process_pdf(file_path, rel_path, metadata, md_splitter, word_splitter):
    file_name = os.path.basename(file_path)
    rel_dir = os.path.dirname(rel_path)
    safe_out_dir = os.path.join(OUT_DIR, rel_dir if rel_dir else "根目录")
    os.makedirs(safe_out_dir, exist_ok=True)
    docs = []

    # 快速探测类型：扫描件（无文字层）只有 MinerU 能处理，后两个引擎帮不上忙
    has_text_layer = _pdf_has_text_layer(file_path)
    if not has_text_layer:
        logger.info(f"    🔍 检测为扫描件，仅 MinerU OCR 可处理")

    md_path = run_mineru(file_path, safe_out_dir)

    if md_path:
        with open(md_path, 'r', encoding='utf-8') as f:
            main_content = f.read()

        if not main_content.strip():
            logger.warning(f"    ⚠️ MinerU 输出为空")
            if not has_text_layer:
                move_to_failed_zone(file_path, file_name, "扫描件 MinerU OCR 输出为空")
                return []
            logger.warning(f"    切换 pymupdf 兜底...")

        elif is_table_garbled(main_content):
            if not has_text_layer:
                # 扫描件乱码：pdfplumber/pymupdf 同样无能为力，直接进失败区
                logger.warning(f"    ⚠️ 扫描件 MinerU 输出乱码，后续引擎无法处理 → 移入失败区")
                move_to_failed_zone(file_path, file_name, "扫描件 MinerU OCR 输出乱码")
                return []
            logger.warning(f"    ⚠️ 检测到乱码 → 切换 pdfplumber")
            table_docs = pdf_table_to_documents(file_path, metadata)
            if table_docs:
                logger.info(f"    ✅ pdfplumber 提取 {len(table_docs)} 块")
                return table_docs
            logger.warning(f"    ⚠️ pdfplumber 无结果，切换 pymupdf 兜底...")

        else:
            # 内嵌图片：调用 Vision LLM 描述，替换 ![](images/xxx) 为 [图示：...]
            main_content = describe_mineru_images(main_content, md_path)
            # 从 content_list.json 注入页码标记（MinerU 本身不在 .md 中写入页码）
            main_content = _inject_page_markers(main_content, md_path)
            # 为每个 chunk 追踪页码：用 null-byte 标记替换页分隔符
            _page_anno_re = re.compile(r'---\s*PAGE\s+(\d+)\s*---', re.IGNORECASE)
            annotated_content = _page_anno_re.sub(
                lambda m: f'\x00PAGE{m.group(1)}\x00', main_content
            )
            raw_chunks = extract_pdf_structured_chunks(annotated_content)
            # 国家规定目录：公式/标准号密集，_has_ocr_noise 误判率高，跳过 OCR 校对
            skip_ocr = rel_path.startswith('国家规定')
            doc_has_noise = False if skip_ocr else _has_ocr_noise(main_content)
            if skip_ocr:
                logger.info(f"    ⏭️  国家规定目录，跳过 OCR 校对，直接入库")
            elif doc_has_noise:
                logger.info(f"    🔧 检测到 OCR 噪音，分批校对（共{len(raw_chunks)}块，每批{OCR_BATCH_SIZE}块，批间暂停{OCR_BATCH_PAUSE}s）...")
            import time as _time
            _page_find_re = re.compile(r'\x00PAGE(\d+)\x00')
            corrected_count, skipped_toc, llm_call_count = 0, 0, 0
            current_page = 1
            for chunk in raw_chunks:
                # 提取页码标注并清除标记
                page_hits = _page_find_re.findall(chunk)
                if page_hits:
                    current_page = int(page_hits[-1])
                clean_chunk = _page_find_re.sub('', chunk).strip()
                if not clean_chunk:
                    continue
                if _is_toc_chunk(clean_chunk):
                    skipped_toc += 1
                    continue
                if doc_has_noise and _has_ocr_noise(clean_chunk):
                    # 每满一批暂停，避免限流
                    if llm_call_count > 0 and llm_call_count % OCR_BATCH_SIZE == 0:
                        logger.info(f"    ⏸️  已校对{llm_call_count}块，暂停{OCR_BATCH_PAUSE}s...")
                        _time.sleep(OCR_BATCH_PAUSE)
                    clean_chunk = llm_correct_ocr(clean_chunk)
                    _time.sleep(OCR_CALL_DELAY)
                    llm_call_count += 1
                    if clean_chunk != chunk:
                        corrected_count += 1
                docs.append(Document(page_content=clean_chunk, metadata={**metadata, 'type': 'pdf', 'page': current_page}))
            if skipped_toc:
                logger.info(f"    🗂️  过滤目录块 {skipped_toc} 个")
            if doc_has_noise:
                logger.info(f"    ✅ LLM 校对完成，调用{llm_call_count}次，修正了 {corrected_count}/{len(raw_chunks)} 块")
            weak_pages = find_weak_pages(md_path)
            if weak_pages:
                logger.warning(f"    🚨 发现 {len(weak_pages)} 个弱页，启动救援...")
                rescue_path = extract_rescue_pdf(file_path, weak_pages)
                if rescue_path:
                    rescue_md = run_mineru(rescue_path, safe_out_dir)
                    if rescue_md:
                        with open(rescue_md, 'r', encoding='utf-8') as f:
                            rescue_content = f.read()
                        if rescue_content.strip():
                            for chunk in extract_pdf_structured_chunks(rescue_content):
                                docs.append(Document(page_content=chunk, metadata={**metadata, 'type': 'pdf_rescue'}))
                            logger.info(f"      ✅ 弱页救援成功")
            return docs
    else:
        logger.warning(f"    ⚠️ MinerU 运行失败")
        if not has_text_layer:
            move_to_failed_zone(file_path, file_name, "扫描件 MinerU OCR 运行失败")
            return []
        logger.warning(f"    切换 pymupdf 兜底...")

    # 只有有文字层的 PDF 才会走到这里
    logger.info(f"    🔄 pymupdf 快速提取中...")
    pymupdf_docs = pymupdf_to_documents(file_path, metadata, word_splitter)
    if pymupdf_docs:
        logger.info(f"    ✅ pymupdf 提取 {len(pymupdf_docs)} 块")
        return pymupdf_docs

    move_to_failed_zone(file_path, file_name, "文字层 PDF 三引擎全部失败")
    return []


# ===================== 4. PPTX 图文处理 =====================

def describe_mineru_images(md_content: str, md_path: str) -> str:
    """扫描 MinerU 输出的 markdown，对每个 ![](images/xxx) 调用 Vision LLM 描述图片内容，
    将标签替换为 [图示：<描述>]，返回替换后的 markdown。"""
    img_dir = os.path.join(os.path.dirname(md_path), "images")
    if not os.path.isdir(img_dir):
        return md_content

    pattern = re.compile(r'!\[\]\(images/([^)]+)\)')
    matches = pattern.findall(md_content)
    if not matches:
        return md_content

    logger.info(f"    🖼️  发现 {len(matches)} 张内嵌图片，逐张调用 Vision LLM 描述...")
    replaced = 0

    def replace_image(m):
        nonlocal replaced
        img_name = m.group(1)
        img_path = os.path.join(img_dir, img_name)
        if not os.path.exists(img_path):
            return m.group(0)
        try:
            ext = os.path.splitext(img_name)[1].lower().lstrip('.')
            mime = 'jpeg' if ext in ('jpg', 'jpeg') else ext
            with open(img_path, 'rb') as f:
                img_b64 = base64.b64encode(f.read()).decode()
            payload = {
                'model': VISION_MODEL,
                'messages': [{'role': 'user', 'content': [
                    {'type': 'image_url', 'image_url': {'url': f'data:image/{mime};base64,{img_b64}'}},
                    {'type': 'text', 'text': (
                        '这是一张工程规范或法律法规文件中的插图。'
                        '请简明描述图示内容：如果有文字标注请转录，如果是示意图/构造图请描述其结构和主要信息，'
                        '如果是表格请转录表格内容。控制在200字以内。'
                    )},
                ]}],
                'temperature': 0.1,
            }
            res = requests.post(
                VISION_BASE_URL,
                headers={'Authorization': f'Bearer {VISION_API_KEY}', 'Content-Type': 'application/json'},
                json=payload, timeout=60
            )
            data = res.json()
            if 'choices' in data:
                desc = data['choices'][0]['message']['content'].strip()
                replaced += 1
                return f'[图示：{desc}]'
        except Exception as e:
            logger.warning(f"    ⚠️  内嵌图片 Vision 失败 ({img_name}): {e}")
        return m.group(0)  # 失败保留原标签

    result = pattern.sub(replace_image, md_content)
    logger.info(f"    ✅ 内嵌图片描述完成：{replaced}/{len(matches)} 张")
    return result


def call_vision_llm_for_slide(img_path: str) -> str:
    """调用 Vision LLM 提取单张幻灯片图片中的全部文字内容。"""
    try:
        with open(img_path, 'rb') as f:
            img_b64 = base64.b64encode(f.read()).decode()
        payload = {
            'model': VISION_MODEL,
            'messages': [{'role': 'user', 'content': [
                {'type': 'image_url', 'image_url': {'url': f'data:image/png;base64,{img_b64}'}},
                {'type': 'text', 'text': '请提取这张PPT幻灯片中的所有文字内容，包括标题、正文、图注、标注、连线旁的文字等，按照原始结构输出。'}
            ]}],
            'temperature': 0.1,
        }
        res = requests.post(
            VISION_BASE_URL,
            headers={'Authorization': f'Bearer {VISION_API_KEY}', 'Content-Type': 'application/json'},
            json=payload, timeout=120
        )
        data = res.json()
        if 'choices' in data:
            return data['choices'][0]['message']['content']
        logger.warning(f"    ⚠️ Vision LLM 返回异常: {data.get('error', data)}")
    except Exception as e:
        logger.warning(f"    ⚠️ Vision LLM 调用失败: {e}")
    return ''


def process_pptx(file_path: str, rel_path: str, metadata: dict) -> list:
    """
    PPTX 处理流程：
    1. LibreOffice 转 PDF
    2. PyMuPDF 渲染每页为 PNG → 存 slide_images/
    3. python-pptx 提取文字：文字充足→直接用；文字不足→调 Vision LLM
    4. 直接返回 Document 列表入库（image_path 存入 metadata，供飞书输出和质检工具使用）
    """
    import fitz
    from pptx import Presentation

    file_name = os.path.basename(file_path)
    base_name = os.path.splitext(file_name)[0]
    safe_base = re.sub(r'[^\w\u4e00-\u9fa5-]', '_', base_name)

    # ── LibreOffice 转 PDF ──
    pdf_conv_dir = os.path.join(OUT_DIR, "_pptx_converted")
    os.makedirs(pdf_conv_dir, exist_ok=True)
    logger.info(f"    🔄 LibreOffice 转换 PPTX → PDF...")
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", pdf_conv_dir, file_path],
            capture_output=True, timeout=300
        )
    except Exception as e:
        logger.error(f"    ❌ LibreOffice 转换失败: {e}")
        return []

    pdf_path = os.path.join(pdf_conv_dir, base_name + ".pdf")
    if not os.path.exists(pdf_path):
        logger.error(f"    ❌ PDF 未生成，跳过: {base_name}")
        return []

    # ── 渲染每页为 PNG ──
    img_subdir = os.path.join(SLIDE_IMG_DIR, safe_base)
    os.makedirs(img_subdir, exist_ok=True)

    # ── 用 python-pptx 预提取文字 ──
    try:
        prs = Presentation(file_path)
        slide_texts = []
        for slide in prs.slides:
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
            slide_texts.append('\n'.join(lines))
    except Exception as e:
        logger.warning(f"    ⚠️ python-pptx 读取失败，全页走 Vision LLM: {e}")
        slide_texts = []

    # ── 逐页渲染 + 提取 → 直接入库 ──
    pdf_doc = fitz.open(pdf_path)
    total   = len(pdf_doc)
    docs = []
    vision_count = 0

    for i, page in enumerate(pdf_doc):
        page_num = i + 1
        img_path = os.path.join(img_subdir, f"p{page_num:04d}.png")
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        pix.save(img_path)

        pptx_text = slide_texts[i] if i < len(slide_texts) else ''

        if len(pptx_text) >= VISION_TEXT_THRESHOLD:
            extracted = pptx_text
        else:
            logger.info(f"    🔍 第{page_num}/{total}页文字不足({len(pptx_text)}字)，调用 Vision LLM...")
            extracted = call_vision_llm_for_slide(img_path)
            vision_count += 1

        if extracted.strip():
            docs.append(Document(
                page_content=extracted,
                metadata={**metadata, 'type': 'pptx', 'page': page_num,
                          'image_path': os.path.relpath(img_path, BASE_DIR)}
            ))

    pdf_doc.close()
    logger.info(f"  ✅ PPTX 处理完成: {total}页，Vision LLM {vision_count}页，生成 {len(docs)} 块")
    return docs


# ===================== 5. PDF 图文处理（非法律法规类）=====================

def process_pdf_with_vision(file_path: str, rel_path: str, metadata: dict) -> list:
    """
    非法律法规 PDF 图文处理流程：
    1. PyMuPDF 渲染每页为 PNG → 存 slide_images/
    2. PyMuPDF 提取文字：文字充足→直接用；文字不足→调 Vision LLM
    3. 直接返回 Document 列表入库（image_path 存入 metadata，供飞书输出和质检工具使用）
    """
    import fitz

    file_name = os.path.basename(file_path)
    base_name = os.path.splitext(file_name)[0]
    safe_base = re.sub(r'[^\w\u4e00-\u9fa5-]', '_', base_name)

    # ── 渲染每页为 PNG 并提取文字 → 直接入库 ──
    img_subdir = os.path.join(SLIDE_IMG_DIR, safe_base)
    os.makedirs(img_subdir, exist_ok=True)

    try:
        pdf_doc = fitz.open(file_path)
    except Exception as e:
        logger.error(f"    ❌ PyMuPDF 打开失败: {e}")
        return []

    total = len(pdf_doc)
    docs = []
    vision_count = 0

    for i, page in enumerate(pdf_doc):
        page_num = i + 1
        img_path = os.path.join(img_subdir, f"p{page_num:04d}.png")
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        pix.save(img_path)

        pymupdf_text = page.get_text().strip()

        if len(pymupdf_text) >= VISION_TEXT_THRESHOLD:
            extracted = pymupdf_text
        else:
            logger.info(f"    🔍 第{page_num}/{total}页文字不足({len(pymupdf_text)}字)，调用 Vision LLM...")
            vision_text = call_vision_llm_for_slide(img_path)
            extracted = vision_text if vision_text else pymupdf_text
            vision_count += 1

        if extracted.strip():
            docs.append(Document(
                page_content=extracted,
                metadata={**metadata, 'type': 'pdf_vision', 'page': page_num,
                          'image_path': os.path.relpath(img_path, BASE_DIR)}
            ))

    pdf_doc.close()
    logger.info(f"  ✅ PDF图文处理完成: {total}页，Vision LLM {vision_count}页，生成 {len(docs)} 块")
    return docs


# ===================== 6. 主程序 =====================
def main():
    import sys
    # --domain EHS案例  只处理指定顶级目录，不传则处理全部
    only_domain = None
    if "--domain" in sys.argv:
        idx = sys.argv.index("--domain")
        only_domain = sys.argv[idx + 1] if idx + 1 < len(sys.argv) else None
    if only_domain:
        logger.info(f"🎯 domain 过滤模式：只处理 [{only_domain}]")
    logger.info("=====================================================")
    logger.info(" 🚀 sync_kb_v11 (EHS案例 Vision LLM 同步)           ")
    logger.info("=====================================================")
    logger.info("🔧 引擎状态检查：")
    logger.info(engine_status())

    logger.info("🔍 正在连接向量数据库...")
    embeddings = HuggingFaceEmbeddings(model_name="BAAI/bge-m3", model_kwargs={'device': 'cpu'})
    db = Chroma(persist_directory=DB_DIR, embedding_function=embeddings)

    md_splitter = MarkdownTextSplitter(chunk_size=800, chunk_overlap=150)
    word_splitter = RecursiveCharacterTextSplitter(
        chunk_size=800, chunk_overlap=150,
        separators=["\n\n", "\n", "。", "！", "？", "；", "，", " "]
    )

    existing_docs = db.get(include=["metadatas"])
    db_files = set([
        meta["source"] for meta in existing_docs.get("metadatas", [])
        if meta and "source" in meta
    ])
    # source → 入库时记录的文件 mtime（旧数据没有该字段，不参与改动检测）
    db_mtimes: dict = {}
    for meta in existing_docs.get("metadatas", []):
        if meta and "source" in meta and meta.get("file_mtime"):
            src = meta["source"]
            if meta["file_mtime"] > db_mtimes.get(src, 0):
                db_mtimes[src] = meta["file_mtime"]

    pdf_files  = glob.glob(os.path.join(KB_DIR, "**/*.pdf"),  recursive=True)
    PDF_files  = glob.glob(os.path.join(KB_DIR, "**/*.PDF"),  recursive=True)
    docx_files = glob.glob(os.path.join(KB_DIR, "**/*.docx"), recursive=True)
    DOCX_files = glob.glob(os.path.join(KB_DIR, "**/*.DOCX"), recursive=True)
    doc_files  = glob.glob(os.path.join(KB_DIR, "**/*.doc"),  recursive=True)
    DOC_files  = glob.glob(os.path.join(KB_DIR, "**/*.DOC"),  recursive=True)
    xlsx_files = glob.glob(os.path.join(KB_DIR, "**/*.xlsx"), recursive=True)
    xls_files  = glob.glob(os.path.join(KB_DIR, "**/*.xls"),  recursive=True)
    pptx_files = glob.glob(os.path.join(KB_DIR, "**/*.pptx"), recursive=True)
    PPTX_files = glob.glob(os.path.join(KB_DIR, "**/*.PPTX"), recursive=True)
    ppt_files  = glob.glob(os.path.join(KB_DIR, "**/*.ppt"),  recursive=True)
    ods_files  = glob.glob(os.path.join(KB_DIR, "**/*.ods"),  recursive=True)
    ODS_files  = glob.glob(os.path.join(KB_DIR, "**/*.ODS"),  recursive=True)

    all_local_files = [
        f for f in (pdf_files + PDF_files + docx_files + DOCX_files + doc_files + DOC_files
                    + xlsx_files + xls_files + pptx_files + PPTX_files + ppt_files
                    + ods_files + ODS_files)
        if '_rescue_' not in f
    ]
    physical_files = set([os.path.relpath(f, KB_DIR) for f in all_local_files])
    rel_to_abs = {os.path.relpath(f, KB_DIR): f for f in all_local_files}

    files_to_delete = db_files - physical_files
    files_to_add    = physical_files - db_files
    # 原地修改检测：入库时记过 mtime 且磁盘文件更新（容差1s）→ 删旧向量+缓存后重入
    files_to_update = {
        rel for rel in (physical_files & db_files)
        if rel in db_mtimes and os.path.getmtime(rel_to_abs[rel]) > db_mtimes[rel] + 1
    }

    if only_domain:
        files_to_add    = {f for f in files_to_add    if Path(f).parts[0] == only_domain}
        files_to_delete = {f for f in files_to_delete if Path(f).parts[0] == only_domain}
        files_to_update = {f for f in files_to_update if Path(f).parts[0] == only_domain}

    logger.info(f"📊 数据库 {len(db_files)} 份 | 本地 {len(physical_files)} 份 | 待删 {len(files_to_delete)} | 待入 {len(files_to_add)} | 待更 {len(files_to_update)}")

    # ── 清理：已删除文件的幽灵数据 + 已修改文件的旧向量（连同 batch_output 解析缓存）──
    if files_to_delete or files_to_update:
        logger.info(f"🗑️ 清理 {len(files_to_delete)} 份幽灵数据、{len(files_to_update)} 份过期数据...")
        for ghost_file in (files_to_delete | files_to_update):
            # 删向量
            db_records = db.get(where={"source": ghost_file})
            if db_records and db_records["ids"]:
                db.delete(ids=db_records["ids"])
                logger.info(f"  🔪 向量已抹除: {ghost_file}")
            # 删 batch_output 解析缓存目录
            ghost_rel_dir  = os.path.dirname(ghost_file)
            ghost_basename = os.path.splitext(os.path.basename(ghost_file))[0]
            cache_dir = os.path.join(OUT_DIR, ghost_rel_dir if ghost_rel_dir else "根目录", ghost_basename)
            if os.path.isdir(cache_dir):
                shutil.rmtree(cache_dir)
                logger.info(f"  🗂️ 解析缓存已清除: {cache_dir}")

    files_to_add = files_to_add | files_to_update

    if not files_to_add:
        logger.info("🎉 知识库已是最新状态，退出。")
        return

    logger.info(f"🌟 发现 {len(files_to_add)} 份新文件，开始处理...")

    for file_path in all_local_files:
        rel_path = os.path.relpath(file_path, KB_DIR)
        if rel_path not in files_to_add:
            continue

        file_name = os.path.basename(file_path)
        timestamp, date_str = extract_time_from_filename(file_name, rel_path)
        rel_dir   = os.path.dirname(rel_path)
        category  = rel_dir if rel_dir else "根目录"

        logger.info(f"  ⏱️ 处理: {file_name}  [{category} | {date_str}]")

        # 顶级目录名作为 domain（EHS案例 / 公司内部 / 国家规定）
        domain = Path(rel_path).parts[0] if Path(rel_path).parts else '未分类'
        metadata = {
            "source":    rel_path,
            "timestamp": timestamp,
            "date_str":  date_str,
            "category":  category,
            "domain":    domain,
            # 入库时的文件 mtime，供下次同步检测"原地修改"（旧数据无此字段则跳过检测）
            "file_mtime": int(os.path.getmtime(file_path)),
        }

        docs_to_insert = []

        # ── Word 文件 ──
        if file_path.lower().endswith(('.docx', '.doc')):
            original_docx_path = file_path  # 保留原始路径供架构提取用
            if file_path.lower().endswith('.doc'):
                logger.info(f"    🔄 .doc → .docx 转换中...")
                converted = convert_doc_to_docx(file_path)
                if converted:
                    file_path = converted
                else:
                    move_to_failed_zone(file_path, file_name, "LibreOffice 转换失败")
                    continue
            try:
                structured_chunks = extract_docx_structured_chunks(file_path)
            except Exception as e:
                logger.warning(f"    ⚠️  结构化切片失败，降级用docx2txt: {e}")
                try:
                    content = docx2txt.process(file_path)
                    structured_chunks = word_splitter.split_text(content) if content.strip() else []
                except Exception as e2:
                    move_to_failed_zone(file_path, file_name, f"Word 提取崩溃: {e2}")
                    continue
            if structured_chunks:
                logger.info(f"    📄 结构化切片: {len(structured_chunks)} 个chunk")
                for seq, chunk in enumerate(structured_chunks, start=1):
                    docs_to_insert.append(Document(
                        page_content=chunk,
                        metadata={**metadata, 'type': 'word', 'page': seq}
                    ))

            else:
                move_to_failed_zone(file_path, file_name, "Word 提取 0 字符（可能是纯图片）")

        # ── Excel / ODS 文件 ──
        elif file_path.lower().endswith(('.xlsx', '.xls', '.ods')):
            if file_path.lower().endswith('.ods'):
                logger.info(f"    🔄 .ods → .xlsx 转换中...")
                converted = convert_ods_to_xlsx(file_path)
                if converted:
                    file_path = converted
                else:
                    move_to_failed_zone(file_path, file_name, "LibreOffice ODS转换失败")
                    continue
            try:
                import openpyxl
                from openpyxl.cell.cell import MergedCell as _MergedCell

                def _cell_val(ws, row, col):
                    """读取单元格值，合并单元格自动追溯到主格。"""
                    cell = ws.cell(row=row, column=col)
                    if isinstance(cell, _MergedCell):
                        for rng in ws.merged_cells.ranges:
                            if cell.coordinate in rng:
                                return ws.cell(rng.min_row, rng.min_col).value
                    return cell.value

                def _chunk_sheet(ws, sheet_name):
                    """
                    合并单元格感知的 Excel 切块：
                    - 自动检测表头行（第一列精确等于常见表头词）
                    - 按第一列（岗位/类别）分组，合并单元格正确追溯
                    - col1 是纯数字序号时，将 col2 描述合并入组标题
                    - 自动过滤签名页脚行
                    - 无表头结构时降级为逐行拼接
                    """
                    max_r, max_c = ws.max_row, ws.max_column
                    rows = []
                    for r in range(1, max_r + 1):
                        row_vals = [str(_cell_val(ws, r, c) or '').strip() for c in range(1, max_c + 1)]
                        rows.append(row_vals)

                    # 精确匹配表头：第一列完全等于常见表头词（避免"项目"误匹配标题行）
                    HEADER_EXACT = {'岗位', '序号', '名称', '类别', '编号', '项目名称'}
                    header_row_idx = None
                    for i, rv in enumerate(rows[:8]):
                        if rv[0] in HEADER_EXACT:
                            header_row_idx = i
                            break

                    def _is_index(s):
                        """判断是否为纯数字序号（如 1、10、2）"""
                        return bool(re.match(r'^\d+$', s)) and len(s) <= 3

                    def _is_footer(rv):
                        """判断是否为签名/页脚行（含冒号占位符，无实质内容）"""
                        joined = ''.join(rv)
                        return joined.count('：') >= 3 and joined.count(' ') > len(joined) * 0.3

                    chunks = []
                    if header_row_idx is not None:
                        data_start = header_row_idx + 1
                        # 跳过与表头第一列相同的子列头行
                        while (data_start < len(rows) and
                               rows[data_start][0] == rows[header_row_idx][0]):
                            data_start += 1

                        current_job, group_lines = None, []

                        def _flush_job(job, lines):
                            if not job or not lines:
                                return
                            content = '\n'.join(ln for ln in lines if ln.strip())
                            if content.strip():
                                chunks.append(f'【{sheet_name} > {job}】\n{content}')

                        for rv in rows[data_start:]:
                            if _is_footer(rv):
                                continue
                            col1 = rv[0]
                            # 跳过合并单元格导致表头词延伸到数据区的行
                            if col1 in HEADER_EXACT:
                                continue
                            col2 = rv[1] if len(rv) > 1 else ''
                            if col1 and col1 != current_job:
                                _flush_job(current_job, group_lines)
                                # 纯数字序号：把 col2 描述并入标题，更具可读性
                                if _is_index(col1) and col2:
                                    current_job = f'{col1}. {col2}'
                                    group_lines = []
                                    # col2 已入标题，正文从 col3 起
                                    line = '  '.join(v for v in rv[2:] if v)
                                else:
                                    current_job = col1
                                    group_lines = []
                                    line = '  '.join(v for v in rv[1:] if v)
                                if line:
                                    group_lines.append(line)
                            else:
                                line = '  '.join(v for v in rv[1:] if v)
                                if line:
                                    group_lines.append(line)
                        _flush_job(current_job, group_lines)
                    else:
                        # 无表头：降级为原始逐行拼接逻辑
                        current_title, current_lines = sheet_name, []
                        for rv in rows:
                            if _is_footer(rv):
                                continue
                            non_empty = [v for v in rv if v]
                            if not non_empty:
                                continue
                            if rv[0] and not any(rv[1:]):
                                content = '\n'.join(ln for ln in current_lines if ln.strip())
                                if content.strip():
                                    chunks.append(f'【{sheet_name} > {current_title}】\n{content}')
                                current_title, current_lines = rv[0], []
                            else:
                                current_lines.append('  '.join(v for v in rv if v))
                        content = '\n'.join(ln for ln in current_lines if ln.strip())
                        if content.strip():
                            chunks.append(f'【{sheet_name} > {current_title}】\n{content}')
                    return chunks

                wb = openpyxl.load_workbook(file_path, data_only=True)
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    for chunk_text in _chunk_sheet(ws, sheet_name):
                        docs_to_insert.append(Document(
                            page_content=chunk_text,
                            metadata={**metadata, 'type': 'excel'}
                        ))
                logger.info(f"    📊 Excel切片: {len(docs_to_insert)} 个chunk")
            except Exception as e:
                logger.warning(f"    ⚠️  Excel 提取失败: {e}")

        # ── PPTX / PPT 文件 ──
        elif file_path.lower().endswith(('.pptx', '.ppt')):
            docs_to_insert = process_pptx(file_path, rel_path, metadata)

        # ── PDF 文件 ──
        elif file_path.lower().endswith('.pdf'):
            # 判断是否走 Vision LLM 路线（EHS案例 或 部门/案例 子目录）
            rel_parts = Path(rel_path).parts
            is_vision_route = (
                domain == 'EHS案例'
                or (len(rel_parts) >= 2 and '案例' in rel_parts[1])
            )
            if is_vision_route:
                # 图文混排，每页渲染 PNG + Vision LLM，带 image_path
                docs_to_insert = process_pdf_with_vision(file_path, rel_path, metadata)
            else:
                # 文字为主，MinerU 结构化提取直接入库
                docs_to_insert = process_pdf(
                    file_path, rel_path, metadata, md_splitter, word_splitter
                )

        # ── 入库（分批避免 CUDA OOM）──
        if docs_to_insert:
            batch_size = 32
            for i in range(0, len(docs_to_insert), batch_size):
                db.add_documents(docs_to_insert[i:i+batch_size])
            logger.info(f"  💾 入库 {len(docs_to_insert)} 块 ✅")
        else:
            logger.warning(f"  ⚠️ {file_name} 未产生任何向量块")

    logger.info("🎉 全部处理完毕。Failed_PDFs 目录存放三引擎均失败的文件。")

if __name__ == "__main__":
    main()
