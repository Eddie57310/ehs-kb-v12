#!/usr/bin/env python3
"""Reindex 设计部/案例 files: direct text extraction + full-page screenshots, no Vision LLM.

Steps:
  1. Extract text from PPTX (python-pptx) and PDF (pymupdf).
  2. Generate full-page screenshots (PPTX via LibreOffice→PDF→pymupdf; PDF via pymupdf directly).
  3. Build chunks: text + [📷 user_crops/{stem}_p{N}_fullpage.png]
  4. Backup chroma_db → delete old chunks → insert new chunks.
  5. Rebuild BM25 index from ChromaDB.
"""

import os, re, shutil, pickle, subprocess, logging, time
from datetime import datetime
from pathlib import Path

import fitz
from pptx import Presentation
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_core.documents import Document

# ── 核心路径 ─────────────────────────────────────────────────────
BASE_DIR    = os.path.expanduser("~/doc_parser_v12")
KB_DIR      = os.path.join(BASE_DIR, "Local_KB")
DB_DIR      = os.path.join(BASE_DIR, "chroma_db")
CROPS_DIR   = os.path.join(BASE_DIR, "user_crops")
BM25_PKL    = os.path.join(BASE_DIR, "bm25_index.pkl")
TMP_PDF_DIR = os.path.join(BASE_DIR, "_tmp_pptx_pdf")

os.makedirs(CROPS_DIR, exist_ok=True)
os.makedirs(TMP_PDF_DIR, exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
)
logger = logging.getLogger(__name__)

# ── 6 个目标文件（v1/v3/用心库v1 已从 PPTX 转为 PDF）──────────────
FILES = [
    "设计部/案例/20241001_惊喜库v1.pdf",
    "设计部/案例/20241101_惊喜库v2.pdf",
    "设计部/案例/20241201_惊喜库v3.pdf",
    "设计部/案例/20241226_杭州公司用心库v1.pdf",
    "设计部/案例/250610_杭州公司用心库v2.pdf",
    "设计部/案例/20260301_华东大区浙江公司沧海桑田计划2.0版.pdf",
]

# 原 PPTX 路径：入库时用的是这些 source，删除旧 chunk 时需一并清理
_LEGACY_PPTX_SOURCES = [
    "设计部/案例/20241001_惊喜库v1.pptx",
    "设计部/案例/20241201_惊喜库v3.pptx",
    "设计部/案例/20241226_杭州公司用心库v1.pptx",
]


# ═══════════════════════════════════════════════════════════════════
#  工具函数
# ═══════════════════════════════════════════════════════════════════

def extract_date(filename: str) -> tuple[int, str]:
    """从文件名提取 (timestamp, date_str)。特殊处理 250610 → 2025-06-10。"""
    # 特殊处理：250610 = 2025-06-10
    if "250610" in filename:
        dt = datetime(2025, 6, 10)
        return int(dt.timestamp()), "2025-06-10"

    m = re.search(r"(\d{4})[-_]?(\d{2})[-_]?(\d{2})", filename)
    if m:
        year, month, day = m.groups()
        date_str = f"{year}-{month}-{day}"
        try:
            dt = datetime.strptime(date_str, "%Y-%m-%d")
            return int(dt.timestamp()), date_str
        except ValueError:
            pass
    return int(datetime(2000, 1, 1).timestamp()), "2000-01-01"


def libreoffice_to_pdf(file_path: str, out_dir: str) -> str | None:
    """LibreOffice 转 PPTX → PDF，返回 pdf 路径或 None（失败时）。"""
    base = os.path.splitext(os.path.basename(file_path))[0]
    pdf_path = os.path.join(out_dir, base + ".pdf")
    if os.path.exists(pdf_path):
        return pdf_path  # 已存在，复用
    logger.info(f"    🔄 LibreOffice 转换中...")
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf",
             "--outdir", out_dir, file_path],
            capture_output=True, timeout=600,
        )
    except subprocess.TimeoutExpired:
        logger.error(f"    ❌ LibreOffice 超时")
        return None
    except Exception as e:
        logger.error(f"    ❌ LibreOffice 异常: {e}")
        return None
    if os.path.exists(pdf_path):
        return pdf_path
    logger.error(f"    ❌ PDF 未生成")
    return None


# ═══════════════════════════════════════════════════════════════════
#  Step 1：文字提取
# ═══════════════════════════════════════════════════════════════════

def extract_pptx_text(file_path: str) -> list[str]:
    """python-pptx 按 shape 顺序遍历，提取每页文字。"""
    prs = Presentation(file_path)
    pages = []
    for slide in prs.slides:
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
        pages.append("\n".join(lines))
    return pages


def extract_pdf_text(file_path: str) -> list[str]:
    """pymupdf 提取 PDF 每页原始文本（去除首尾空行）。"""
    doc = fitz.open(file_path)
    pages = []
    for page in doc:
        text = page.get_text("text").strip()
        pages.append(text)
    doc.close()
    return pages


# ═══════════════════════════════════════════════════════════════════
#  Step 2：整页截图
# ═══════════════════════════════════════════════════════════════════

def screenshot_pdf(pdf_path: str, stem: str, total: int) -> list[str]:
    """用 pymupdf 渲染 PDF 每页为 fullpage PNG，返回相对路径列表。"""
    doc = fitz.open(pdf_path)
    mat = fitz.Matrix(2.0, 2.0)  # ~144 DPI
    rel_paths = []
    for i, page in enumerate(doc):
        pn = i + 1
        fname = f"{stem}_p{pn}_fullpage.png"
        img_path = os.path.join(CROPS_DIR, fname)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        pix.save(img_path)
        rel_paths.append(f"user_crops/{fname}")
    doc.close()
    return rel_paths


# ═══════════════════════════════════════════════════════════════════
#  Step 4-5：数据库操作
# ═══════════════════════════════════════════════════════════════════

def backup_chroma():
    """备份当前 chroma_db。"""
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    bak = os.path.join(BASE_DIR, f"chroma_db.bak_{ts}_设计案例重制前")
    logger.info(f"📦 备份 chroma_db → {os.path.basename(bak)} ...")
    shutil.copytree(DB_DIR, bak)
    logger.info(f"✅ 备份完成")


def rebuild_bm25():
    """从 ChromaDB 读取全部文档，序列化到 bm25_index.pkl（供 feishu_ws_server 热加载）。"""
    import chromadb
    logger.info("🔨 重建 BM25 索引...")
    c = chromadb.PersistentClient(path=DB_DIR)
    col = c.get_collection("langchain")
    r = col.get(include=["documents", "metadatas"])
    corpus = list(zip(r["documents"], r["metadatas"]))
    with open(BM25_PKL + ".tmp", "wb") as f:
        pickle.dump(corpus, f)
    os.replace(BM25_PKL + ".tmp", BM25_PKL)
    logger.info(f"✅ BM25 重建完成：{len(corpus)} 条 → {os.path.basename(BM25_PKL)}")


# ═══════════════════════════════════════════════════════════════════
#  Main
# ═══════════════════════════════════════════════════════════════════

def main():
    logger.info("=" * 60)
    logger.info(" 设计部/案例 数据重制 - 直接文本 + 整页截图")
    logger.info("=" * 60)

    # ── 初始化 ChromaDB ──
    logger.info("🔍 连接向量数据库...")
    embeddings = HuggingFaceEmbeddings(
        model_name="BAAI/bge-m3",
        model_kwargs={"device": "cpu"},
    )
    db = Chroma(persist_directory=DB_DIR, embedding_function=embeddings)

    # ── Step 4.1：备份 ──
    backup_chroma()

    all_new_docs = []

    for rel_path in FILES:
        file_path = os.path.join(KB_DIR, rel_path)
        if not os.path.exists(file_path):
            logger.warning(f"⚠️  文件不存在，跳过: {rel_path}")
            continue

        file_name = os.path.basename(file_path)
        stem = os.path.splitext(file_name)[0]
        ext = os.path.splitext(file_name)[1].lower()
        timestamp, date_str = extract_date(file_name)
        domain = Path(rel_path).parts[0]
        category = str(Path(rel_path).parent)

        logger.info(f"📄 处理: {file_name}  [{category} | {date_str}]")

        # ── 文字提取 ──
        if ext == ".pptx":
            page_texts = extract_pptx_text(file_path)
        elif ext == ".pdf":
            page_texts = extract_pdf_text(file_path)
        else:
            logger.warning(f"    ⚠️  不支持格式: {ext}")
            continue

        total = len(page_texts)
        logger.info(f"    共 {total} 页")

        # ── 截图 ──
        if ext == ".pptx":
            # PPTX → LibreOffice → PDF → pymupdf 渲染
            tmp_pdf = libreoffice_to_pdf(file_path, TMP_PDF_DIR)
            if not tmp_pdf:
                logger.error(f"    ❌ LibreOffice 转换失败，跳过: {file_name}")
                continue
            img_rel_paths = screenshot_pdf(tmp_pdf, stem, total)
            if len(img_rel_paths) != total:
                logger.warning(f"    ⚠️  页数不一致：pptx={total}，截图={len(img_rel_paths)}，图文可能错位！")
            # 清理临时 PDF
            try:
                os.remove(tmp_pdf)
            except Exception:
                pass
        else:
            # PDF 直接渲染
            img_rel_paths = screenshot_pdf(file_path, stem, total)

        # ── Step 3：构建 chunk ──
        for i in range(total):
            text = page_texts[i].strip()
            if i < len(img_rel_paths):
                img_tag = f"[📷 {img_rel_paths[i]}]"
            else:
                img_tag = "[📷 缺图]"
            chunk_text = text + f"\n\n{img_tag}" if text else img_tag

            meta = {
                "domain":   domain,
                "source":   rel_path,
                "date_str": date_str,
                "timestamp": timestamp,
                "category": category,
                "page":     i + 1,
                "type":     ext.lstrip("."),
            }
            all_new_docs.append(Document(page_content=chunk_text, metadata=meta))

        logger.info(f"    ✅ 生成 {total} 个 chunk")

    if not all_new_docs:
        logger.error("❌ 没有生成任何 chunk，终止")
        return

    # ── Step 4.2：删除旧 chunk（按 source，含旧 PPTX 路径）──
    sources_to_delete = FILES + _LEGACY_PPTX_SOURCES
    logger.info(f"🗑️  清理旧 chunk（{len(sources_to_delete)} 个 source）...")
    for rel_path in sources_to_delete:
        records = db.get(where={"source": rel_path})
        if records and records["ids"]:
            db.delete(ids=records["ids"])
            logger.info(f"    🔪 已删除: {rel_path}  ({len(records['ids'])} 块)")

    # ── Step 4.3：入库新 chunk ──
    logger.info(f"💾 入库 {len(all_new_docs)} 个新 chunk...")
    batch_size = 32
    for i in range(0, len(all_new_docs), batch_size):
        db.add_documents(all_new_docs[i:i + batch_size])
        logger.info(f"    批次 {i // batch_size + 1}: {len(all_new_docs[i:i + batch_size])} 块")
    logger.info("✅ 入库完成")

    # ── Step 5：重建 BM25 ──
    rebuild_bm25()

    # ── 清理临时目录 ──
    try:
        shutil.rmtree(TMP_PDF_DIR)
    except Exception:
        pass

    logger.info("=" * 60)
    logger.info(" 🎉 设计部/案例 数据重制完成！")
    logger.info("=" * 60)


if __name__ == "__main__":
    main()
